# Biblioteca para trabalhar com arquivos JSON
import json

# Biblioteca para ler arquivos Excel
import pandas as pd

# Biblioteca para manipular caminhos de arquivos
from pathlib import Path

# Biblioteca para ler arquivos PowerPoint
from pptx import Presentation

# Estrutura usada pelo LangChain para armazenar textos
from langchain_core.documents import Document

# Leitores prontos para diferentes tipos de arquivos
from langchain_community.document_loaders import (
    PyPDFLoader,
    TextLoader,
    CSVLoader,
    UnstructuredHTMLLoader,
    UnstructuredWordDocumentLoader,
)


# ============================================================
# FUNÇÕES PARA LEITURA DOS ARQUIVOS
# ============================================================

def ler_pdf(caminho):
    """Lê um arquivo PDF."""
    leitor = PyPDFLoader(str(caminho))
    return leitor.load()


def ler_word(caminho):
    """Lê um documento do Word."""
    leitor = UnstructuredWordDocumentLoader(str(caminho))
    return leitor.load()


def ler_texto(caminho):
    """Lê arquivos TXT e Markdown."""
    leitor = TextLoader(
        str(caminho),
        encoding="utf-8"
    )
    return leitor.load()


def ler_html(caminho):
    """Lê arquivos HTML."""
    leitor = UnstructuredHTMLLoader(str(caminho))
    return leitor.load()


def ler_csv(caminho):
    """Lê arquivos CSV."""
    leitor = CSVLoader(
        str(caminho),
        encoding="utf-8"
    )
    return leitor.load()


def ler_excel(caminho):
    """
    Um arquivo Excel pode possuir várias abas.

    Esta função lê todas elas e transforma
    cada aba em um Document.
    """

    documentos = []

    abas = pd.read_excel(
        caminho,
        sheet_name=None
    )

    for nome_aba in abas:

        tabela = abas[nome_aba]

        texto = tabela.to_string(index=False)

        documento = Document(
            page_content=texto,
            metadata={
                "aba": nome_aba
            }
        )

        documentos.append(documento)

    return documentos


def ler_powerpoint(caminho):
    """Lê arquivos PowerPoint."""

    documentos = []

    apresentacao = Presentation(str(caminho))

    numero_slide = 0

    for slide in apresentacao.slides:

        numero_slide += 1

        texto_slide = ""

        for forma in slide.shapes:

            if hasattr(forma, "text"):

                texto_slide += forma.text + "\n"

        if texto_slide.strip():

            documento = Document(
                page_content=texto_slide,
                metadata={
                    "slide": numero_slide
                }
            )

            documentos.append(documento)

    return documentos


def ler_json(caminho):
    """Lê arquivos JSON."""

    with open(caminho, encoding="utf-8") as arquivo:

        dados = json.load(arquivo)

    texto = json.dumps(
        dados,
        ensure_ascii=False,
        indent=2
    )

    return [
        Document(
            page_content=texto,
            metadata={}
        )
    ]


# ============================================================
# FUNÇÃO PRINCIPAL
# ============================================================

def ler_arquivo(caminho):

    """
    Recebe o caminho de qualquer arquivo
    e escolhe automaticamente qual função
    deve ser utilizada para fazer a leitura.
    """

    extensao = Path(caminho).suffix.lower()

    if extensao == ".pdf":
        documentos = ler_pdf(caminho)

    elif extensao == ".docx":
        documentos = ler_word(caminho)

    elif extensao in [".txt", ".md"]:
        documentos = ler_texto(caminho)

    elif extensao in [".html", ".htm"]:
        documentos = ler_html(caminho)

    elif extensao == ".csv":
        documentos = ler_csv(caminho)

    elif extensao == ".xlsx":
        documentos = ler_excel(caminho)

    elif extensao == ".pptx":
        documentos = ler_powerpoint(caminho)

    elif extensao == ".json":
        documentos = ler_json(caminho)

    else:
        documentos = []

    # Guarda o nome do arquivo nos metadados
    for documento in documentos:
        documento.metadata["arquivo"] = Path(caminho).name

    return documentos